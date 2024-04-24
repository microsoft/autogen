import markdownify
import pptx

from autogen.coding.func_with_reqs import with_requirements


@with_requirements(["python-pptx", "beautifulsoup4", "markdownify"], ["pptx", "markdownify"])
def pptx_to_md(local_path):
    """
    Convert a PowerPoint presentation (PPTX) to Markdown format.

    Args:
        local_path (str): The local path to the PPTX file.

    Returns:
        str: The Markdown content generated from the PPTX file.
    """
    # A simplified version of https://github.com/microsoft/autogen/blob/266cefc1737e0077667bce441c541a90865582b1/autogen/browser_utils/mdconvert.py
    # Will change the code into importing the MdConverter Class as soon as PR#1929 is merged.
    import html
    import re

    from bs4 import BeautifulSoup

    def _is_picture(shape):
        # Check if shape is a picture
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "image"):
                return True
        return False

    def _is_table(shape):
        # Check if shape is a table
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
            return True
        return False

    def _convert(html_content):
        """Helper function that converts and HTML string."""

        # Parse the string
        soup = BeautifulSoup(html_content, "html.parser")

        # Remove javascript and style blocks
        for script in soup(["script", "style"]):
            script.extract()

        # Print only the main content
        body_elm = soup.find("body")
        webpage_text = ""
        if body_elm:
            webpage_text = markdownify.MarkdownConverter().convert_soup(body_elm)
        else:
            webpage_text = markdownify.MarkdownConverter().convert_soup(soup)
        return webpage_text

    md_content = ""
    presentation = pptx.Presentation(local_path)
    slide_num = 0
    for slide in presentation.slides:
        slide_num += 1

        md_content += f"\n\n<!-- Slide number: {slide_num} -->\n"

        title = slide.shapes.title
        for shape in slide.shapes:
            # Pictures
            if _is_picture(shape):
                alt_text = ""
                try:
                    alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
                except Exception:
                    pass

                filename = re.sub(r"\W", "", shape.name) + ".jpg"
                md_content += "\n![" + (alt_text if alt_text else shape.name) + "](" + filename + ")\n"

            # Tables
            if _is_table(shape):
                html_table = "<html><body><table>"
                first_row = True
                for row in shape.table.rows:
                    html_table += "<tr>"
                    for cell in row.cells:
                        if first_row:
                            html_table += "<th>" + html.escape(cell.text) + "</th>"
                        else:
                            html_table += "<td>" + html.escape(cell.text) + "</td>"
                    html_table += "</tr>"
                    first_row = False
                html_table += "</table></body></html>"
                md_content += "\n" + _convert(html_table).text_content.strip() + "\n"

            # Text areas
            elif shape.has_text_frame:
                if shape == title:
                    md_content += "# " + shape.text.lstrip() + " "
                else:
                    md_content += shape.text + " "

    md_content = md_content.strip()

    if slide.has_notes_slide:
        md_content += "\n\n### Notes:\n"
        notes_frame = slide.notes_slide.notes_text_frame
        if notes_frame is not None:
            md_content += notes_frame.text
        md_content = md_content.strip()

    return md_content
